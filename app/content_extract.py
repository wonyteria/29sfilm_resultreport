from __future__ import annotations

import csv
from dataclasses import asdict, dataclass
from pathlib import Path


TEXT_LIMIT = 12000


@dataclass
class EvidenceFile:
    path: str
    extension: str
    text: str
    extraction: str
    error: str = ""


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    return "\n".join((page.extract_text() or "") for page in PdfReader(path).pages)


def _read_docx(path: Path) -> str:
    from docx import Document

    document = Document(path)
    chunks = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(chunks)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    chunks: list[str] = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text)
    return "\n".join(chunks)


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    chunks: list[str] = []
    try:
        for ws in wb.worksheets:
            chunks.append(f"[{ws.title}]")
            for row in ws.iter_rows(values_only=True):
                values = [str(value) for value in row if value not in (None, "")]
                if values:
                    chunks.append(" | ".join(values))
                if sum(len(x) for x in chunks) >= TEXT_LIMIT:
                    return "\n".join(chunks)
    finally:
        wb.close()
    return "\n".join(chunks)


def _read_csv(path: Path) -> str:
    for encoding in ("utf-8-sig", "cp949", "utf-8"):
        try:
            with path.open("r", encoding=encoding, newline="") as handle:
                return "\n".join(" | ".join(row) for row in csv.reader(handle))
        except UnicodeDecodeError:
            continue
    return ""


READERS = {".pdf": _read_pdf, ".docx": _read_docx, ".pptx": _read_pptx, ".xlsx": _read_xlsx, ".csv": _read_csv}


def extract_project_text(root_value: str, relative_paths: list[str]) -> list[EvidenceFile]:
    root = Path(root_value)
    evidence: list[EvidenceFile] = []
    for relative in relative_paths:
        path = root / relative
        reader = READERS.get(path.suffix.lower())
        if not reader:
            continue
        try:
            text = reader(path)[:TEXT_LIMIT].strip()
            evidence.append(EvidenceFile(relative, path.suffix.lower(), text, "본문 추출" if text else "본문 없음"))
        except Exception as exc:
            evidence.append(EvidenceFile(relative, path.suffix.lower(), "", "추출 실패", str(exc)))
    return evidence


def evidence_as_dict(items: list[EvidenceFile]) -> list[dict]:
    return [asdict(item) for item in items]
