from __future__ import annotations

import json
import re
from collections import Counter
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation


@dataclass(frozen=True)
class PageModule:
    key: str
    label: str
    section: str
    mode: str
    patterns: tuple[str, ...]
    source_hints: tuple[str, ...] = ()


MODULES = [
    PageModule("cover", "표지", "공통", "required", ("결과보고서",)),
    PageModule("index", "목차", "공통", "required", ("index", "contents")),
    PageModule("section_overview", "사업개요 섹션", "사업개요", "required", ("01 사업개요",)),
    PageModule("overview", "사업개요", "사업개요", "required", ("1-1 개요",), ("최종기획안", "사업계획서")),
    PageModule("section_promotion", "사업홍보 섹션", "사업홍보", "required", ("02 사업홍보",)),
    PageModule("promotion_schedule", "홍보 일정", "사업홍보", "conditional", ("사업홍보 2-1 일정", "2-1 일정"), ("최종기획안", "일정표")),
    PageModule("poster_video", "포스터·홍보영상", "사업홍보", "conditional", ("포스터 디자인", "포스터 및 홍보영상", "홍보영상"), ("포스터", "홍보영상")),
    PageModule("poster_distribution", "포스터 배포", "사업홍보", "repeat", ("포스터 배포",), ("포스터배포", "배포목록")),
    PageModule("news", "기사·보도자료", "사업홍보", "repeat", ("보도자료", "개최 기사", "시상식 기사", "온라인 기사", "지면 기사"), ("보도자료", "온라인기사", "지면기사")),
    PageModule("print_ad", "지면광고", "사업홍보", "repeat", ("5단광고", "지면 광고", "전면광고"), ("5단광고", "신문광고")),
    PageModule("homepage", "홈페이지", "사업홍보", "conditional", ("홈페이지",), ("홈페이지" ,)),
    PageModule("outdoor_ad", "옥외·교통광고", "사업홍보", "repeat", ("버스", "gbus", "옥외", "지하철", "스크린도어"), ("버스광고", "gbus", "교통광고", "옥외광고")),
    PageModule("message", "카카오톡·SMS·뉴스레터", "사업홍보", "repeat", ("카카오톡", "sms", "뉴스레터", "스티비"), ("카카오톡", "sms", "뉴스레터", "스티비")),
    PageModule("community", "공모전 사이트·커뮤니티", "사업홍보", "repeat", ("커뮤니티", "올콘", "공모전 사이트"), ("올콘", "커뮤니티", "공모전사이트")),
    PageModule("sns", "SNS 운영", "사업홍보", "repeat", ("sns 운영", "인스타그램", "페이스북", "유튜브 게시물"), ("인스타", "페이스북", "유튜브")),
    PageModule("digital_ad", "디지털 광고", "사업홍보", "repeat", ("디지털 광고", "구글", "메타", "네이버 광고", "카카오 광고"), ("구글", "카카오", "메타", "네이버", "광고")),
    PageModule("supporters", "서포터즈·홍보대사", "사업홍보", "repeat", ("서포터즈", "홍보대사", "제작지원"), ("서포터즈", "홍보대사")),
    PageModule("section_submission", "출품·심사 섹션", "출품·심사", "required", ("03 공모 및 심사", "03 출품결과")),
    PageModule("submission", "출품 결과", "출품·심사", "required", ("출품 결과", "출품 현황", "출품작 수"), ("출품현황", "출품목록")),
    PageModule("judging", "심사", "출품·심사", "conditional", ("3-2 심사", "심사위원", "심사 방식"), ("심사기획안", "심사위원")),
    PageModule("award_list", "수상작 목록", "수상작", "required", ("수상작 no", "수상작 목록", "상격 감독명"), ("수상작목록",)),
    PageModule("award_detail", "수상작 상세", "수상작", "repeat", ("시놉시스",), ("대표이미지", "시놉시스", "영상url")),
    PageModule("section_ceremony", "시상식 섹션", "시상식", "conditional", ("04 시상식",)),
    PageModule("ceremony_overview", "시상식 개요·큐시트", "시상식", "conditional", ("4-1 개요", "시상식 개요", "큐시트", "식순"), ("시상식기획안", "큐시트")),
    PageModule("ceremony_production", "시상식 제작물", "시상식", "repeat", ("제작물 운영", "상장", "상패", "포토월", "배너"), ("제작물",)),
    PageModule("ceremony_photos", "시상식 현장", "시상식", "repeat", ("현장 및 스케치", "행사 전경", "행사사진", "현장사진"), ("현장사진",)),
    PageModule("reviews", "후기", "성과·총평", "repeat", ("영화제 후기", "인스타그램 후기", "블로그 후기", "관객 후기"), ("후기",)),
    PageModule("summary", "총평", "성과·총평", "required", ("총평",), ("통합성과", "총평메모")),
    PageModule("closing", "마무리", "공통", "required", ("감사합니다", "고맙습니다")),
]


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", text.replace("|", " ")).strip().lower()


def slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            chunks.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                chunks.extend(cell.text for cell in row.cells if cell.text)
    return normalize(" ".join(chunks))


def classify_slide(text: str, slide_number: int, total: int) -> PageModule | None:
    if slide_number == 1:
        return MODULES[0]
    if slide_number == 2 and any(x in text for x in ("index", "contents")):
        return next(x for x in MODULES if x.key == "index")
    if slide_number == total and any(x in text for x in ("감사합니다", "고맙습니다")):
        return next(x for x in MODULES if x.key == "closing")
    if len(text) < 80:
        section_rules = (
            ("01 사업개요", "section_overview"),
            ("02 사업홍보", "section_promotion"),
            ("03 공모 및 심사", "section_submission"),
            ("03 출품결과", "section_submission"),
            ("04 시상식", "section_ceremony"),
        )
        for marker, key in section_rules:
            if marker in text:
                return next(x for x in MODULES if x.key == key)
    scores: list[tuple[int, PageModule]] = []
    for module in MODULES:
        if module.key in {"cover", "index", "section_overview", "section_promotion", "section_submission", "section_ceremony", "closing"}:
            continue
        hits = sum(1 for pattern in module.patterns if pattern.lower() in text)
        if hits:
            specificity = max(len(pattern) for pattern in module.patterns if pattern.lower() in text)
            scores.append((hits * 100 + specificity, module))
    return max(scores, key=lambda item: item[0])[1] if scores else None


def analyze_presentation(path: str | Path) -> dict:
    path = Path(path)
    presentation = Presentation(path)
    pages = []
    counts: Counter[str] = Counter()
    for index, slide in enumerate(presentation.slides, start=1):
        text = slide_text(slide)
        module = classify_slide(text, index, len(presentation.slides))
        key = module.key if module else "unclassified"
        counts[key] += 1
        pages.append({
            "slide": index,
            "module": key,
            "label": module.label if module else "미분류",
            "text_preview": text[:240],
            "shapes": len(slide.shapes),
        })
    return {"file": path.name, "slides": len(presentation.slides), "module_counts": dict(counts), "pages": pages}


def build_catalog(paths: Iterable[str | Path], output: str | Path) -> dict:
    reports = [analyze_presentation(path) for path in paths]
    payload = {"schema_version": 1, "modules": [asdict(module) for module in MODULES], "reports": reports}
    Path(output).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return payload
